import os
import re
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def clean_text(text):
    # 마크다운 제거
    text = text.replace('**', '').replace('__', '')
    # 단순 영어 번역 괄호 제거 (예: (Correlation), (Prediction) 등)
    # 한글이 포함된 괄호는 보존
    text = re.sub(r'\s*\([A-Za-z\s,./-]+\)', '', text)
    return text.strip()

def create_ppt_from_md(md_path, output_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    # '## 슬라이드 N: 제목' 기준으로 분할
    slides_raw = re.split(r'## 슬라이드 \d+:', content)[1:]
    
    for slide_text in slides_raw:
        lines = slide_text.strip().split('\n')
        if not lines: continue
        
        # 첫 줄은 제목
        title = clean_text(lines[0].strip())
        
        # 불렛 포인트 파싱
        bullets = []
        for line in lines[1:]:
            line = line.strip()
            if line.startswith('- '):
                bullets.append(clean_text(line[2:]))
            elif line.startswith('• '):
                bullets.append(clean_text(line[2:]))
                
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 헤더 배경
        header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(245, 246, 250)
        header.line.fill.background()

        # 제목 텍스트 박스
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(40, 53, 147)
        
        # 1차시 0번 슬라이드 (강의 개요) 예외 처리
        if "강의 개요" in title:
            for i, text in enumerate(bullets):
                y_pos = 2.5 + (i * 1.2)
                txt_box = slide.shapes.add_textbox(Inches(2.0), Inches(y_pos), Inches(9.33), Inches(1.0))
                p = txt_box.text_frame.paragraphs[0]
                p.text = f"• {text}"
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = RGBColor(30, 30, 30)
            continue
            
        # 비주얼 가이드 박스 (좌측)
        frame = slide.shapes.add_shape(1, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(240, 240, 240)
        frame.line.color.rgb = RGBColor(3, 169, 244)
        frame.line.width = Pt(1)
        
        guide = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.4), Inches(4.1))
        tf = guide.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "[VISUAL GUIDE]\n추후 관련 이미지 또는 그래프 삽입"
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER
        
        # 본문 불렛 포인트 (우측) - 첫/끝 위치 절대 고정
        FIX_TOP = 1.8
        FIX_BOTTOM = 6.0
        
        num_s = len(bullets)
        gap = (FIX_BOTTOM - FIX_TOP) / (num_s - 1) if num_s > 1 else 0
        
        # 동적 폰트 조절 (세션이 4개 이내이므로 30~32pt 유지 가능)
        best_fs = 32 if num_s <= 3 else 28
        if num_s > 4: best_fs = 24
        
        for i, text in enumerate(bullets):
            y_pos = FIX_TOP + (i * gap)
            txt_box = slide.shapes.add_textbox(Inches(6.8), Inches(y_pos), Inches(6.0), Inches(1.2))
            tf = txt_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {text}"
            p.font.size = Pt(best_fs)
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(35, 35, 35)

    prs.save(output_path)

if __name__ == "__main__":
    base_dir = "/home/tj/ml_presentation/0304/"
    # REVISED.md 파일들 찾기
    md_files = sorted(glob.glob(os.path.join(base_dir, "*_REVISED.md")))
    for md_path in md_files:
        out_name = os.path.basename(md_path).replace('_REVISED.md', '_FINAL.pptx')
        out_path = os.path.join(base_dir, out_name)
        create_ppt_from_md(md_path, out_path)
        print(f"[생성 완료] {out_path}")
